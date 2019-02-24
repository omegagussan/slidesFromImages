
# coding: utf-8
#science
import cv2
import numpy as np
import imutils

#python
from os import listdir
from os import mkdir
from os.path import isfile, join
import argparse
import sys, os
import shutil

#pptx
from pptx import Presentation

#helpers
def order_points(pts):
    rect = np.zeros((4, 2), dtype = "float32")
    s = pts.sum(axis = 1)
    rect[0] = pts[np.argmin(s)]
    rect[2] = pts[np.argmax(s)]
    diff = np.diff(pts, axis = 1)
    rect[1] = pts[np.argmin(diff)]
    rect[3] = pts[np.argmax(diff)]
    return rect

def four_point_transform(image, pts):
    rect = order_points(pts)
    (tl, tr, br, bl) = rect
    widthA = np.sqrt(((br[0] - bl[0]) ** 2) + ((br[1] - bl[1]) ** 2))
    widthB = np.sqrt(((tr[0] - tl[0]) ** 2) + ((tr[1] - tl[1]) ** 2))
    maxWidth = max(int(widthA), int(widthB))
    heightA = np.sqrt(((tr[0] - br[0]) ** 2) + ((tr[1] - br[1]) ** 2))
    heightB = np.sqrt(((tl[0] - bl[0]) ** 2) + ((tl[1] - bl[1]) ** 2))
    maxHeight = max(int(heightA), int(heightB))
    dst = np.array([
        [0, 0],
        [maxWidth - 1, 0],
        [maxWidth - 1, maxHeight - 1],
        [0, maxHeight - 1]], dtype = "float32")
    M = cv2.getPerspectiveTransform(rect, dst)
    return cv2.warpPerspective(image, M, (maxWidth, maxHeight))

SUPPORTED_FORMATS = (".BMP", ".DIB", ".JPEG", ".JPG", ".JPE", ".JP2", ".PNG", ".WEBP", ".PBM", ".PGM", ".PPM", ".SR", ".RAS", ".TIFF", ".TIF") #tuple
def isImage(filename):
    return filename.upper().endswith(SUPPORTED_FORMATS)

# construct the argument parser and parse the arguments
parser = argparse.ArgumentParser()
parser.add_argument("-i", "--input_folder", required = True,
	help = "Folder with images to be parsed")
parser.add_argument("-o", "--output_file", required = True,
	help = "Name of slides to be produced")
parser.add_argument('--r', dest='feature', action='store_true')
parser.set_defaults(feature=False)
args = vars(parser.parse_args())

input_folder = args["input_folder"]
output_file = args["output_file"]
rotate_input = (args["feature"])
temp_folder = "./temp"
os.mkdir(temp_folder)

only_files = [f for f in listdir(input_folder) if isfile(join(input_folder, f))]
only_files.sort()
only_images = list(filter(lambda x: isImage(x), only_files))

TARGET_SIZE = 500
images = []
edges = []
originals = []

for image in only_images:
    relative_path = input_folder + "/" + image
    print("processing: " + relative_path)
    image = cv2.imread(relative_path)
    if rotate_input:
        original=cv2.transpose(image)
        image=cv2.transpose(image)

    else:
        original = image.copy()
    ratio = image.shape[0] / TARGET_SIZE
    image = imutils.resize(image, height = TARGET_SIZE)
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    blur = cv2.GaussianBlur(image,(5,5),0)
    edge = cv2.Canny(blur, 75, 200)
    kernel = np.ones((5,5),np.uint8)
    edge = cv2.dilate(edge,kernel,iterations = 1)    
    images.append(image)
    edges.append(edge)
    originals.append(original)
    
print("small images: " + str(len(images)))
print("color images: " + str(len(originals)))
print("edge images: " + str(len(edges)))

cnt_images = []
for idx in range(len(edges)):
    cnts = cv2.findContours(edges[idx].copy(), cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
    cnts = imutils.grab_contours(cnts)
    cnts = sorted(cnts, key = cv2.contourArea, reverse = True)[:5]
    for cnt in cnts:
        peri = cv2.arcLength(cnt, True)
        approx = cv2.approxPolyDP(cnt, 0.05 * peri, True)
        if len(approx) == 4:
            screenCnt = approx
            cnt_images.append(screenCnt)
            break

for idx in range(len(originals)):
    warped = four_point_transform(originals[idx], cnt_images[idx].reshape(4, 2) * ratio)
    cv2.imwrite(temp_folder + "/" + only_images[idx] , warped)

prs = Presentation()
BLANK_SLIDE = prs.slide_layouts[6]
slideAspectRatio = prs.slide_width/prs.slide_height
for image in only_images:
    relative_path = temp_folder + "/" + image
    slide = prs.slides.add_slide(BLANK_SLIDE)
    shapes = slide.shapes
    picture = shapes.add_picture(relative_path, 0, 0, prs.slide_width, prs.slide_height)
    
prs.save(output_file)
shutil.rmtree(temp_folder)